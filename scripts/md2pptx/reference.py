# Copyright (c) 2025 TrailLensCo
# All rights reserved.
#
# This file is proprietary and confidential.

"""Reference document generation for background images."""

import logging
from pathlib import Path
from typing import Optional

from pptx import Presentation
from pptx.util import Inches

logger = logging.getLogger(__name__)


class ReferenceDocGenerator:
    """Generator for creating reference PPTX documents with background images."""

    @staticmethod
    def create_with_background(
        background_image: Path,
        output_path: Path,
        aspect_ratio: str = "16:9",
    ) -> Path:
        """
        Create a reference PPTX with a background image on all slides.

        Args:
            background_image: Path to background image file.
            output_path: Path to save reference PPTX.
            aspect_ratio: Presentation aspect ratio ("16:9" or "4:3").

        Returns:
            Path to created reference PPTX.

        Raises:
            FileNotFoundError: If background image doesn't exist.
            ValueError: If image format is unsupported.
        """
        if not background_image.exists():
            raise FileNotFoundError(f"Background image not found: {background_image}")

        # Create presentation
        prs = Presentation()

        # Set slide dimensions based on aspect ratio
        if aspect_ratio == "16:9":
            prs.slide_width = Inches(10)
            prs.slide_height = Inches(5.625)
        elif aspect_ratio == "4:3":
            prs.slide_width = Inches(10)
            prs.slide_height = Inches(7.5)
        else:
            logger.warning(f"Unknown aspect ratio '{aspect_ratio}', using 16:9")
            prs.slide_width = Inches(10)
            prs.slide_height = Inches(5.625)

        # Get the blank slide layout (usually index 6, but may vary)
        blank_layout = None
        for layout in prs.slide_layouts:
            if layout.name.lower() in ("blank", "blank layout"):
                blank_layout = layout
                break

        # Fallback to first layout if blank not found
        if blank_layout is None:
            blank_layout = prs.slide_layouts[0]
            logger.warning("Blank layout not found, using first available layout")

        # Create a master slide with background
        slide = prs.slides.add_slide(blank_layout)

        # Add background image to fill the slide
        try:
            slide.shapes.add_picture(
                str(background_image),
                0,  # left
                0,  # top
                prs.slide_width,
                prs.slide_height,
            )
            logger.info(f"Added background image: {background_image}")
        except Exception as e:
            raise ValueError(f"Failed to add background image: {e}")

        # Save reference document
        prs.save(str(output_path))
        logger.info(f"Created reference PPTX with background: {output_path}")

        return output_path

    @staticmethod
    def validate_image(image_path: Path) -> bool:
        """
        Validate that the image file is in a supported format.

        Args:
            image_path: Path to image file.

        Returns:
            True if valid, False otherwise.
        """
        if not image_path.exists():
            logger.error(f"Image file not found: {image_path}")
            return False

        supported_extensions = {".jpg", ".jpeg", ".png", ".gif", ".bmp", ".svg"}
        if image_path.suffix.lower() not in supported_extensions:
            logger.error(
                f"Unsupported image format: {image_path.suffix}. "
                f"Supported: {', '.join(supported_extensions)}"
            )
            return False

        return True
